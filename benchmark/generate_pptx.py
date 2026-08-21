#!/usr/bin/env python3
"""Build benchmark/RESULTS.pptx from the scored xCodeEval JSONL data."""
import os, json
from collections import defaultdict, Counter
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

BASE = os.path.dirname(os.path.abspath(__file__))
OUT  = os.path.join(BASE, "RESULTS.pptx")

MODELS = [
    ("claude-opus-5",   "Claude Opus 5"),
    ("claude-sonnet-5", "Claude Sonnet 5"),
    ("gpt-5.6-sol",     "GPT-5.6-sol"),
    ("gpt4o",           "GPT-5.1 (gpt4o)"),
    ("laguna-nvfp4",    "Laguna NV-FP4"),
    ("nemotron-550b",   "Nemotron-550B"),
    ("qwen-nvfp4",      "Qwen NV-FP4"),
]
LANGS = [("C++","GNU C++17"),("Go","Go"),("Java","Java 17"),
         ("Javascript","Node.js"),("Kotlin","Kotlin 1.4"),("PHP","PHP"),("Python","PyPy 3")]
LANG_NAMES = [l for l,_ in LANGS]

# ---- palette ----
NAVY   = RGBColor(0x0F,0x2A,0x43)
BLUE   = RGBColor(0x1F,0x6F,0xB2)
LIGHT  = RGBColor(0xE8,0xF0,0xF7)
WHITE  = RGBColor(0xFF,0xFF,0xFF)
DARK   = RGBColor(0x22,0x2B,0x33)
GREY   = RGBColor(0x6B,0x75,0x80)
GOOD   = RGBColor(0x1B,0x7F,0x4B)   # green text for best
ZEBRA  = RGBColor(0xF4,0xF7,0xFA)

# per-category tints for the failure-breakdown columns (pale headers)
FAIL_TINT = {
    "WRONG_ANSWER":          RGBColor(0xFC,0xEF,0xC7),  # pale yellow
    "RUNTIME_ERROR":         RGBColor(0xF9,0xD6,0xD2),  # pale red
    "COMPILATION_ERROR":     RGBColor(0xDD,0xD9,0xF0),  # pale purple
    "TIME_LIMIT_EXCEEDED":   RGBColor(0xFC,0xE0,0xC4),  # pale orange
    "MEMORY_LIMIT_EXCEEDED": RGBColor(0xDD,0xE6,0xEC),  # pale grey-blue
}
# failure categories, in fixed column order, with short header labels
CATS  = ["WRONG_ANSWER","RUNTIME_ERROR","COMPILATION_ERROR",
         "TIME_LIMIT_EXCEEDED","MEMORY_LIMIT_EXCEEDED"]
CSHORT = {"WRONG_ANSWER":"Wrong","RUNTIME_ERROR":"Runtime",
          "COMPILATION_ERROR":"Compile","TIME_LIMIT_EXCEEDED":"TimeLim",
          "MEMORY_LIMIT_EXCEEDED":"MemLim"}
PERFECT = RGBColor(0x1B,0x7F,0x4B)  # green text for "all passed"

def ps_path(m,c): return os.path.join(BASE,m,"ps","reproduce_1",f"{c}.jsonl")
def ct_path(m,c): return os.path.join(BASE,m,"ct_compact_small",
                     "eval_code_translation_compact_small_execeval",f"{c}.jsonl")

def estimator(n,c,k):
    if n-c<k: return 1.0
    p=1.0
    for i in range(n-c+1,n+1): p*=(1.0-k/i)
    return 1.0-p

def load(path):
    if not os.path.exists(path) or os.path.getsize(path)==0: return None
    res=defaultdict(list); fail=Counter(); nsamp=[]
    with open(path) as f:
        for line in f:
            line=line.strip()
            if not line: continue
            s=json.loads(line)
            uid=s["source_data"]["src_uid"]
            attempts=s["unit_test_results"]; nsamp.append(len(attempts))
            for a in attempts:
                if isinstance(a,dict) and "error" in a: continue
                res[uid].append(a)
                first=None; ok=True
                for t in a:
                    if t["exec_outcome"]!="PASSED":
                        ok=False
                        if first is None: first=t["exec_outcome"]
                if not ok: fail[first or "UNKNOWN"]+=1
    if not res: return None
    total=[]; correct=[]
    for v in res.values():
        passed=[all(t["exec_outcome"]=="PASSED" for t in a) for a in v]
        total.append(len(passed)); correct.append(sum(passed))
    def pk(k):
        if not all(t>=k for t in total): return None
        return 100.0*sum(estimator(t,c,k) for t,c in zip(total,correct))/len(total)
    return {"p1":pk(1),"p5":pk(5),"fail":fail,"nmin":min(nsamp),"nmax":max(nsamp)}

DATA={}
for m,_ in MODELS:
    for lang,comp in LANGS:
        DATA[("PS",m,lang)]=load(ps_path(m,comp))
        DATA[("CT",m,lang)]=load(ct_path(m,comp))

ISSUES = {
"claude-opus-5":[
 ("Sampling limit","Azure Anthropic API supports only n=1 → pass@1 exact, pass@5 not computable."),
 ("Run","PS 560 files, CT 1688 files; both finished clean (exit=0)."),
 ("Result","Strongest model overall. Weak spot: Javascript CT 71.6% (40 RUNTIME_ERRORs)."),
],
"claude-sonnet-5":[
 ("Sampling limit","Same Azure n=1 constraint → pass@1 only, pass@5 blank."),
 ("JS collapse","PS Javascript 32.3% — 36 RUNTIME_ERRORs, a systematic model failure."),
 ("Result","Solid elsewhere; PHP CT weak at 68.0%."),
],
"gpt-5.6-sol":[
 ("Samples","n=5 → pass@1 and pass@5 both available."),
 ("Rerun needed","First eval left gaps; PS/CT reruns filled missing compiler outputs."),
 ("Kotlin broken","Catastrophic COMPILATION_ERROR (PS 180 / CT 572) → Kotlin CT 49.8%."),
 ("Result","Best sampling model — CT avg 90.2% @5."),
],
"gpt4o":[
 ("Alias","'gpt4o' endpoint routes to GPT-5.1 via LiteLLM/Azure."),
 ("Sampling cap","Azure enforces n≤8 (not 20). pass@1/pass@5 both valid."),
 ("Scope","CT ran the FULL compact split (400 problems/lang, 2804 files, ~2h)."),
 ("Result","Weak on Kotlin (485 compile errors) and PHP (503 runtime errors)."),
],
"laguna-nvfp4":[
 ("Backend","Local vLLM NV-FP4, enable_thinking=False, max_tokens=4096."),
 ("Samples","PS n=20, CT n=8; some CT langs n=0..8 (empty/short generations)."),
 ("Weakness","Heavy COMPILATION_ERROR on compiled langs (Go 289 PS, Kotlin 471 PS)."),
 ("Result","Lowest PS of the FP4 pair (41.9% @1)."),
],
"nemotron-550b":[
 ("Baseline subset","Fixed baseline idx/lang pairs (PS 418, CT 1723) — comparable to gpt4o."),
 ("PS n=1","PS is n=1 for all langs except C++ (partial n=5) → PS pass@5 blank."),
 ("CT n=5","CT is clean n=5 → CT pass@5 shown."),
 ("Weakness","Dominant COMPILATION_ERROR on CT Go/Kotlin/Java; PHP floor 43.8%."),
],
"qwen-nvfp4":[
 ("Backend","Local vLLM NV-FP4; port moved 5000→5001 mid-run, reruns needed."),
 ("Samples","PS n=20, CT n=8."),
 ("Missing data","PS Kotlin: no usable output. PS Python: only 4 scorable problems"),
 ("","(16.2% is on a tiny sample, not comparable)."),
 ("Compile-heavy","Worst compiler reliability: Go PS 818 / Go CT 466 COMPILATION_ERRORs."),
],
}

# ---------------- slide helpers ----------------
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height

def add_slide():
    return prs.slides.add_slide(BLANK)

def band(slide, title, subtitle=None):
    bar = slide.shapes.add_shape(1, 0, 0, SW, Inches(1.15))
    bar.fill.solid(); bar.fill.fore_color.rgb = NAVY; bar.line.fill.background()
    tf = bar.text_frame; tf.word_wrap=True
    tf.margin_left=Inches(0.5); tf.margin_top=Inches(0.12)
    p=tf.paragraphs[0]; p.text=title
    p.font.size=Pt(28); p.font.bold=True; p.font.color.rgb=WHITE
    if subtitle:
        p2=tf.add_paragraph(); p2.text=subtitle
        p2.font.size=Pt(13); p2.font.color.rgb=RGBColor(0xB9,0xD2,0xE8)

def textbox(slide,left,top,width,height):
    tb=slide.shapes.add_textbox(left,top,width,height)
    tb.text_frame.word_wrap=True
    return tb.text_frame

def style_cell(cell, text, *, bold=False, size=13, color=DARK, fill=None,
               align=PP_ALIGN.CENTER):
    cell.text=str(text)
    cell.vertical_anchor=MSO_ANCHOR.MIDDLE
    cell.margin_left=Inches(0.04); cell.margin_right=Inches(0.04)
    cell.margin_top=Inches(0.02); cell.margin_bottom=Inches(0.02)
    if fill is not None:
        cell.fill.solid(); cell.fill.fore_color.rgb=fill
    else:
        cell.fill.solid(); cell.fill.fore_color.rgb=WHITE
    for p in cell.text_frame.paragraphs:
        p.alignment=align
        for r in p.runs:
            r.font.size=Pt(size); r.font.bold=bold; r.font.color.rgb=color

def metric_table_slide(task, metric, title, sub):
    slide=add_slide(); band(slide,title,sub)
    rows=len(MODELS)+1; cols=len(LANG_NAMES)+2
    left=Inches(0.35); top=Inches(1.45)
    width=SW-Inches(0.7); height=Inches(0.62)*rows
    gtbl=slide.shapes.add_table(rows,cols,left,top,width,height)
    tbl=gtbl.table
    tbl.columns[0].width=Inches(2.6)
    for i in range(1,cols): tbl.columns[i].width=Inches((13.333-0.7-2.6)/(cols-1))
    # header
    style_cell(tbl.cell(0,0),"Model",bold=True,size=14,color=WHITE,fill=BLUE,align=PP_ALIGN.LEFT)
    for j,l in enumerate(LANG_NAMES):
        style_cell(tbl.cell(0,j+1),l,bold=True,size=12,color=WHITE,fill=BLUE)
    style_cell(tbl.cell(0,cols-1),"Avg",bold=True,size=13,color=WHITE,fill=NAVY)
    # best per column
    colvals={j:[] for j in range(len(LANG_NAMES))}
    for m,_ in MODELS:
        for j,l in enumerate(LANG_NAMES):
            r=DATA[(task,m,l)]
            if r and r[metric] is not None: colvals[j].append(r[metric])
    colmax={j:(max(v) if v else None) for j,v in colvals.items()}
    for i,(m,disp) in enumerate(MODELS):
        z = ZEBRA if i%2 else WHITE
        style_cell(tbl.cell(i+1,0),disp,bold=True,size=12,color=DARK,fill=z,align=PP_ALIGN.LEFT)
        vals=[]
        for j,l in enumerate(LANG_NAMES):
            r=DATA[(task,m,l)]
            v=r[metric] if r else None
            if v is None:
                style_cell(tbl.cell(i+1,j+1),"—",size=12,color=GREY,fill=z)
            else:
                best = colmax[j] is not None and abs(v-colmax[j])<1e-9
                style_cell(tbl.cell(i+1,j+1),f"{v:.1f}%",bold=best,size=12,
                           color=GOOD if best else DARK, fill=z)
                vals.append(v)
        avg=f"{sum(vals)/len(vals):.1f}%" if vals else "—"
        style_cell(tbl.cell(i+1,cols-1),avg,bold=True,size=12,color=NAVY,fill=LIGHT)
    # footnote
    tf=textbox(slide,Inches(0.35),SH-Inches(0.55),SW-Inches(0.7),Inches(0.45))
    p=tf.paragraphs[0]
    p.text="Green = best in column. — = not computable (n<5) or no data. Values are % (unbiased pass@k estimator)."
    p.font.size=Pt(11); p.font.color.rgb=GREY

def title_slide():
    slide=add_slide()
    bg=slide.shapes.add_shape(1,0,0,SW,SH)
    bg.fill.solid(); bg.fill.fore_color.rgb=NAVY; bg.line.fill.background()
    tf=bg.text_frame; tf.word_wrap=True
    tf.margin_left=Inches(0.9); tf.margin_top=Inches(2.2)
    p=tf.paragraphs[0]; p.text="xCodeEval Benchmark"
    p.font.size=Pt(48); p.font.bold=True; p.font.color.rgb=WHITE
    for line,sz,col in [
        ("pass@1 & pass@5  —  7 models × 7 languages",22,RGBColor(0xCF,0xE2,0xF3)),
        ("Program Synthesis (PS)  •  Code Translation (CT)",18,RGBColor(0x9F,0xC0,0xDD)),
        ("Generated 2026-08-21",14,RGBColor(0x7F,0xA5,0xC8)),
    ]:
        q=tf.add_paragraph(); q.text=line; q.font.size=Pt(sz); q.font.color.rgb=col

def legend_slide():
    slide=add_slide(); band(slide,"How to read these results","Metrics, sampling, and why some cells are blank")
    tf=textbox(slide,Inches(0.6),Inches(1.5),SW-Inches(1.2),SH-Inches(1.9))
    items=[
        ("pass@1","Unbiased estimator; equals correct/total when one sample. Reported for every model.",True),
        ("pass@5","Unbiased estimator 1 − C(n−c,5)/C(n,5). Needs n≥5 samples for EVERY problem in a language, else blank.",True),
        ("—","Not computable for that model+language.",True),
        ("",""  ,False),
        ("Blank pass@5 rows — why:","",True),
        ("Claude Opus 5 & Sonnet 5","Azure Anthropic API supports only n=1 → no pass@5 at all.",False),
        ("Nemotron-550B (PS)","PS generated at n=1 (C++ partial n=5) → PS pass@5 blank; CT is n=5 so CT pass@5 shown.",False),
        ("Qwen NV-FP4 (Kotlin PS)","No usable generations → blank in every metric.",False),
    ]
    first=True
    for k,v,hd in items:
        p=tf.paragraphs[0] if first else tf.add_paragraph(); first=False
        if k and v:
            r=p.add_run(); r.text=f"{k}  —  "; r.font.bold=True; r.font.size=Pt(16); r.font.color.rgb=BLUE
            r2=p.add_run(); r2.text=v; r2.font.size=Pt(15); r2.font.color.rgb=DARK
        elif k:
            r=p.add_run(); r.text=k; r.font.bold=True; r.font.size=Pt(17); r.font.color.rgb=NAVY
        p.space_after=Pt(8)
    # failure-category colour swatches (used on per-model slides)
    sy=SH-Inches(1.15)
    lab=textbox(slide,Inches(0.6),sy-Inches(0.35),SW-Inches(1.2),Inches(0.3))
    lp=lab.paragraphs[0]; lp.text="Failure-mode colours (per-model slides):"
    lp.font.size=Pt(13); lp.font.bold=True; lp.font.color.rgb=NAVY
    x=Inches(0.6)
    for c in CATS:
        sw=slide.shapes.add_shape(1,x,sy,Inches(0.3),Inches(0.3))
        sw.fill.solid(); sw.fill.fore_color.rgb=FAIL_TINT[c]
        sw.line.color.rgb=GREY
        tb=slide.shapes.add_textbox(x+Inches(0.34),sy-Inches(0.02),Inches(2.0),Inches(0.35))
        tp=tb.text_frame.paragraphs[0]; tp.text=CSHORT[c]
        tp.font.size=Pt(12); tp.font.color.rgb=DARK
        x=x+Inches(2.35)

COST_PS=[
    ("Claude Opus 5","1","390.1K / 758.1K","$5.00 / $25.00","$20.90"),
    ("Claude Sonnet 5","1","384.9K / 1.51M","$2.00 / $10.00","$15.91"),
    ("GPT-5.1 (dep: gpt4o)","8","272.5K / 1.93M","$2.00 / $8.00","$15.96"),
    ("GPT-5.6-sol","5","272.5K / 3.25M","$5.00 / $20.00","$66.36"),
    ("Nemotron-550B","5","215.9K / 6.34M","— / —","—"),
    ("Laguna NV-FP4","20","303.6K / 6.52M","— / —","—"),
    ("Qwen NV-FP4","20","213.9K / 12.79M","— / —","—"),
]
COST_CT=[
    ("Claude Opus 5","1","936.0K / 1.09M","$5.00 / $25.00","$31.94"),
    ("Claude Sonnet 5","1","1.12M / 1.29M","$2.00 / $10.00","$15.10"),
    ("GPT-5.1 (dep: gpt4o)","8","1.27M / 8.30M","$2.00 / $8.00","$68.93"),
    ("GPT-5.6-sol","5","754.7K / 3.31M","$5.00 / $20.00","$70.04"),
    ("Nemotron-550B","5","825.4K / 11.33M","— / —","—"),
    ("Laguna NV-FP4","8","879.0K / 4.77M","— / —","—"),
    ("Qwen NV-FP4","8","809.3K / 5.11M","— / —","—"),
]

def cost_table_slide(task_label, sub, data, total_line):
    slide=add_slide(); band(slide,task_label,sub)
    header=("Model","n","Input / Output Tokens","Rate /1M (in/out)","Cost")
    rows=len(data)+1; cols=5
    gtbl=slide.shapes.add_table(rows,cols,Inches(1.15),Inches(1.5),Inches(11.0),Inches(0.55)*rows)
    tbl=gtbl.table
    widths=[3.1,0.8,3.1,2.6,1.4]
    for i,w in enumerate(widths): tbl.columns[i].width=Inches(w)
    for j,h in enumerate(header):
        style_cell(tbl.cell(0,j),h,bold=True,size=13,color=WHITE,fill=BLUE,
                   align=PP_ALIGN.LEFT if j==0 else PP_ALIGN.CENTER)
    for i,row in enumerate(data,1):
        z=ZEBRA if i%2 else WHITE
        local = row[4]=="—"
        for j,val in enumerate(row):
            col = GREY if (local and j>=3) else DARK
            style_cell(tbl.cell(i,j),val,size=12,color=col,fill=z,
                       bold=(j==4 and not local),
                       align=PP_ALIGN.LEFT if j==0 else PP_ALIGN.CENTER)
    tf=textbox(slide,Inches(1.15),Inches(1.5)+Inches(0.55)*rows+Inches(0.15),Inches(11.0),Inches(0.8))
    p=tf.paragraphs[0]; p.text=total_line
    p.font.size=Pt(15); p.font.bold=True; p.font.color.rgb=NAVY
    p2=tf.add_paragraph()
    p2.text="Local NV-FP4 models (Nemotron, Laguna, Qwen) run on owned GPUs — no per-token $ rate."
    p2.font.size=Pt(11); p2.font.color.rgb=GREY

def task_block(slide, m, task, label, top):
    """One full-width table: Language | pass@1 | pass@5 | <5 failure-share cols>."""
    cap=textbox(slide,Inches(0.4),top,Inches(12.5),Inches(0.3))
    cp=cap.paragraphs[0]; cp.text=label
    cp.font.size=Pt(14); cp.font.bold=True; cp.font.color.rgb=BLUE
    top=top+Inches(0.32)
    cols=3+len(CATS)                      # Language, @1, @5, + categories
    rows=1+len(LANG_NAMES)
    rowh=Inches(0.34)
    gtbl=slide.shapes.add_table(rows,cols,Inches(0.4),top,Inches(12.5),rowh*rows)
    tbl=gtbl.table
    widths=[1.9,1.25,1.25]+[1.56]*len(CATS)
    for i,w in enumerate(widths): tbl.columns[i].width=Inches(w)
    hdr=["Language","pass@1","pass@5"]+[CSHORT[c] for c in CATS]
    for j,h in enumerate(hdr):
        fill = FAIL_TINT[CATS[j-3]] if j>=3 else BLUE
        col  = DARK if j>=3 else WHITE
        style_cell(tbl.cell(0,j),h,bold=True,size=11,color=col,fill=fill,
                   align=PP_ALIGN.LEFT if j==0 else PP_ALIGN.CENTER)
    for i,lang in enumerate(LANG_NAMES,1):
        z=ZEBRA if i%2 else WHITE
        style_cell(tbl.cell(i,0),lang,bold=True,size=11,color=DARK,fill=z,align=PP_ALIGN.LEFT)
        r=DATA[(task,m,lang)]
        if not r:
            for j in range(1,cols):
                style_cell(tbl.cell(i,j),"—",size=11,color=GREY,fill=z)
            continue
        p1="—" if r["p1"] is None else f'{r["p1"]:.1f}%'
        p5="—" if r["p5"] is None else f'{r["p5"]:.1f}%'
        style_cell(tbl.cell(i,1),p1,size=11,color=DARK,fill=z)
        style_cell(tbl.cell(i,2),p5,size=11,color=GREY if p5=="—" else DARK,fill=z)
        tot=sum(r["fail"].values())
        for j,c in enumerate(CATS,3):
            if tot==0:
                txt,cc=("clean" if j==3 else ""),PERFECT
            else:
                share=100*r["fail"].get(c,0)/tot
                txt = f"{share:.0f}%" if share>0 else "·"
                cc  = DARK if share>0 else GREY
            style_cell(tbl.cell(i,j),txt,size=11,color=cc,fill=z)
    return top+rowh*rows

def model_slide(m,disp):
    slide=add_slide()
    band(slide,disp,"Per-language pass@k  +  full failure-mode breakdown "
                    "(each category = % share of that language's failed attempts)")
    y=Inches(1.35)
    y=task_block(slide,m,"PS","Program Synthesis (PS)",y)+Inches(0.18)
    y=task_block(slide,m,"CT","Code Translation (CT)",y)
    tf=textbox(slide,Inches(0.4),SH-Inches(0.42),Inches(12.5),Inches(0.35))
    p=tf.paragraphs[0]
    p.text=("Failure columns share a colour with the legend.  "
            "· = 0% of failures.  “clean” = language had no failed attempts.  "
            "— = no data.")
    p.font.size=Pt(10); p.font.color.rgb=GREY

# ---- build deck ----
title_slide()
legend_slide()
metric_table_slide("PS","p1","Table 1 — Program Synthesis · pass@1","Fraction of problems solved by a single sample (%)")
metric_table_slide("PS","p5","Table 2 — Program Synthesis · pass@5","At least one of 5 samples solves the problem (%)")
metric_table_slide("CT","p1","Table 3 — Code Translation · pass@1","Single-sample translation correctness (%)")
metric_table_slide("CT","p5","Table 4 — Code Translation · pass@5","Best-of-5 translation correctness (%)")
cost_table_slide("Tokens & Cost — Program Synthesis (PS)",
                 "Per-model token usage and API cost",
                 COST_PS, "PS total (API models only): $119.14")
cost_table_slide("Tokens & Cost — Code Translation (CT)",
                 "Per-model token usage and API cost   •   Grand total PS+CT: $305.15",
                 COST_CT, "CT total (API models only): $186.01")
for m,disp in MODELS:
    model_slide(m,disp)

prs.save(OUT)
print("wrote", OUT, "slides:", len(prs.slides._sldIdLst))

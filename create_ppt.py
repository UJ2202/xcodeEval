from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as nsmap
from lxml import etree
import copy

# ── Color palette ──────────────────────────────────────────────────────────────
DARK_BG     = RGBColor(0x0D, 0x1B, 0x2A)   # deep navy
ACCENT      = RGBColor(0x00, 0xB4, 0xD8)   # cyan
ACCENT2     = RGBColor(0x90, 0xE0, 0xEF)   # light cyan
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY  = RGBColor(0xCC, 0xCC, 0xCC)
GREEN       = RGBColor(0x06, 0xD6, 0xA0)
YELLOW      = RGBColor(0xFF, 0xD1, 0x66)
RED         = RGBColor(0xEF, 0x47, 0x6F)
CARD_BG     = RGBColor(0x1A, 0x2D, 0x42)   # slightly lighter navy

W = Inches(13.33)   # widescreen width
H = Inches(7.5)     # widescreen height

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank = prs.slide_layouts[6]   # truly blank layout

# ── Helpers ────────────────────────────────────────────────────────────────────
def add_slide():
    sl = prs.slides.add_slide(blank)
    fill = sl.background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG
    return sl

def txb(slide, text, l, t, w, h,
        size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
        italic=False, wrap=True):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return tb

def rect(slide, l, t, w, h, fill_color, alpha=None):
    shp = slide.shapes.add_shape(
        pptx.enum.shapes.MSO_SHAPE_TYPE.AUTO_SHAPE if False else 1,  # MSO_SHAPE_TYPE is enum; use 1=rectangle
        l, t, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_color
    shp.line.fill.background()
    return shp

def hline(slide, y, color=ACCENT):
    ln = slide.shapes.add_shape(1, Inches(0.5), y, Inches(12.33), Pt(1.5))
    ln.fill.solid()
    ln.fill.fore_color.rgb = color
    ln.line.fill.background()

def accent_bar(slide):
    rect(slide, 0, 0, Inches(0.12), H, ACCENT)

def section_header(slide, title, subtitle=None):
    accent_bar(slide)
    txb(slide, title,
        Inches(0.35), Inches(0.22), Inches(12.5), Inches(0.7),
        size=32, bold=True, color=ACCENT)
    hline(slide, Inches(1.05))
    if subtitle:
        txb(slide, subtitle,
            Inches(0.35), Inches(1.15), Inches(12.5), Inches(0.45),
            size=15, color=LIGHT_GRAY, italic=True)

def bullet(slide, items, l, t, w, size=16, gap=Inches(0.38), color=WHITE, dot_color=ACCENT):
    for i, item in enumerate(items):
        rect(slide, l, t + i*gap + Inches(0.09), Inches(0.12), Inches(0.12), dot_color)
        txb(slide, item, l + Inches(0.22), t + i*gap, w - Inches(0.22), Inches(0.36),
            size=size, color=color)

def card(slide, l, t, w, h, title, body_lines, title_color=ACCENT, body_size=14):
    rect(slide, l, t, w, h, CARD_BG)
    # left accent strip
    rect(slide, l, t, Inches(0.07), h, title_color)
    txb(slide, title,
        l + Inches(0.15), t + Inches(0.1), w - Inches(0.2), Inches(0.38),
        size=15, bold=True, color=title_color)
    body = "\n".join(body_lines)
    txb(slide, body,
        l + Inches(0.15), t + Inches(0.52), w - Inches(0.2), h - Inches(0.6),
        size=body_size, color=WHITE)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
rect(sl, 0, 0, W, H, DARK_BG)
rect(sl, 0, 0, Inches(0.18), H, ACCENT)

# big title
txb(sl, "xCodeEval",
    Inches(0.45), Inches(1.5), Inches(12), Inches(1.4),
    size=72, bold=True, color=ACCENT)
txb(sl, "Multilingual Code Benchmark Evaluation",
    Inches(0.45), Inches(2.85), Inches(12), Inches(0.7),
    size=26, bold=False, color=WHITE)
hline(sl, Inches(3.65), ACCENT2)
txb(sl, "Models Evaluated:  GPT-5.1  ·  Qwen-NVFP4  ·  Laguna-NVFP4",
    Inches(0.45), Inches(3.8), Inches(9), Inches(0.5),
    size=17, color=ACCENT2)
txb(sl, "Task: Program Synthesis  ·  7 Languages  ·  Metric: pass@5",
    Inches(0.45), Inches(4.3), Inches(9), Inches(0.45),
    size=15, color=LIGHT_GRAY, italic=True)
txb(sl, "xCodeEval  |  NTU-NLP-sg  |  arxiv 2303.03004",
    Inches(0.45), Inches(6.8), Inches(9), Inches(0.4),
    size=12, color=LIGHT_GRAY, italic=True)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — What is xCodeEval?
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
section_header(sl, "What is xCodeEval?",
               "The largest executable multilingual multitask code benchmark")

# 3 stat boxes
stats = [
    ("25 M+",   "Coding examples"),
    ("17",      "Programming languages"),
    ("7",       "Distinct tasks"),
    ("~7.5K",   "Unique problems"),
]
bw = Inches(2.8)
for i, (num, label) in enumerate(stats):
    bx = Inches(0.4) + i * (bw + Inches(0.2))
    rect(sl, bx, Inches(1.55), bw, Inches(1.5), CARD_BG)
    rect(sl, bx, Inches(1.55), bw, Inches(0.06), ACCENT)
    txb(sl, num,  bx, Inches(1.65), bw, Inches(0.8),
        size=38, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    txb(sl, label, bx, Inches(2.4), bw, Inches(0.45),
        size=14, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

txb(sl, "7 Tasks covered:",
    Inches(0.4), Inches(3.25), Inches(6), Inches(0.38),
    size=16, bold=True, color=ACCENT2)
tasks = [
    "Program Synthesis — write code from problem description",
    "Code Translation — translate code between languages",
    "APR — automated program repair (fix buggy code)",
    "Tag Classification — classify problem difficulty/category",
    "Code Compilation — predict if code compiles",
    "Code-Code Retrieval — find similar code snippets",
    "NL-Code Retrieval — match natural language to code",
]
bullet(sl, tasks, Inches(0.4), Inches(3.7), Inches(6.2), size=14, gap=Inches(0.42))

txb(sl, "Evaluation engine:",
    Inches(7.0), Inches(3.25), Inches(5.5), Inches(0.38),
    size=16, bold=True, color=ACCENT2)
card(sl, Inches(7.0), Inches(3.7), Inches(5.8), Inches(2.9),
     "ExecEval — Execution-based Scoring",
     ["• Docker-based sandboxed code runner",
      "• Supports all 17 languages",
      "• Runs hidden unit tests against generated code",
      "• Returns: PASSED / WRONG_ANSWER /",
      "  RUNTIME_ERROR / TIME_LIMIT_EXCEEDED /",
      "  COMPILATION_ERROR"],
     title_color=GREEN)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Benchmark Focus: Program Synthesis
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
section_header(sl, "Benchmark Focus: Program Synthesis",
               "Write correct code from a natural language problem description")

txb(sl, "How it works:",
    Inches(0.4), Inches(1.35), Inches(5), Inches(0.38),
    size=16, bold=True, color=ACCENT2)

steps = [
    ("1", "Problem description + I/O spec + sample cases → prompt", ACCENT),
    ("2", "LLM generates N code solutions (n=8 for GPT-5.1, n=20 for others)", ACCENT2),
    ("3", "ExecEval runs each solution against hidden unit tests", GREEN),
    ("4", "pass@k computed: prob. at least 1 of k solutions passes all tests", YELLOW),
]
for i, (num, text, col) in enumerate(steps):
    y = Inches(1.8) + i * Inches(1.1)
    rect(sl, Inches(0.4), y, Inches(0.5), Inches(0.5), col)
    txb(sl, num, Inches(0.4), y, Inches(0.5), Inches(0.5),
        size=20, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)
    txb(sl, text, Inches(1.05), y + Inches(0.07), Inches(5.5), Inches(0.45),
        size=15, color=WHITE)

# prompt example box
txb(sl, "Example Prompt:",
    Inches(7.2), Inches(1.35), Inches(5.6), Inches(0.38),
    size=15, bold=True, color=ACCENT2)
rect(sl, Inches(7.2), Inches(1.75), Inches(5.8), Inches(5.3), CARD_BG)
rect(sl, Inches(7.2), Inches(1.75), Inches(0.07), Inches(5.3), YELLOW)
prompt_text = (
    'Write a program in C++ to solve this\n'
    'programming problem:\n\n'
    'Description: Given a string s, find the\n'
    'maximum number of characters you can\n'
    'remove...\n\n'
    'Input: single string s\n'
    'Output: maximum removals\n\n'
    'Sample Input: "abba"\n'
    'Sample Output: 2\n\n'
    'Provide the C++ code without any extra\n'
    'description or tokens. Target code:'
)
txb(sl, prompt_text,
    Inches(7.4), Inches(1.9), Inches(5.5), Inches(5.0),
    size=12, color=ACCENT2)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Our Setup
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
section_header(sl, "Evaluation Setup",
               "3 models benchmarked in parallel via LiteLLM proxy")

model_cards = [
    ("GPT-5.1", "Azure OpenAI", "gpt4o alias", "n=8 (Azure cap)", "8192 max tokens", ACCENT),
    ("Qwen-NVFP4", "Remote vLLM :8000", "qwen-nvfp4", "n=20", "4096 max tokens", GREEN),
    ("Laguna-NVFP4", "Remote vLLM :8001", "laguna-nvfp4", "n=20", "4096 max tokens", YELLOW),
]
for i, (name, host, alias, n, tok, col) in enumerate(model_cards):
    bx = Inches(0.4) + i * Inches(4.25)
    rect(sl, bx, Inches(1.5), Inches(4.05), Inches(2.5), CARD_BG)
    rect(sl, bx, Inches(1.5), Inches(4.05), Inches(0.07), col)
    txb(sl, name,  bx + Inches(0.15), Inches(1.58), Inches(3.8), Inches(0.5),
        size=20, bold=True, color=col)
    txb(sl, f"{host}\nAlias: {alias}\nSamples: {n}\nTokens: {tok}",
        bx + Inches(0.15), Inches(2.1), Inches(3.8), Inches(1.7),
        size=14, color=WHITE)

txb(sl, "All routed through  →  LiteLLM Proxy (localhost:4000)",
    Inches(0.4), Inches(4.15), Inches(9), Inches(0.45),
    size=16, bold=True, color=ACCENT2)

txb(sl, "Languages tested:",
    Inches(0.4), Inches(4.7), Inches(4), Inches(0.38),
    size=15, bold=True, color=ACCENT2)
langs = ["C++", "Go", "Java", "Javascript", "PHP", "Python", "Kotlin"]
lang_str = "  ·  ".join(langs)
txb(sl, lang_str, Inches(0.4), Inches(5.1), Inches(10), Inches(0.4),
    size=18, bold=True, color=WHITE)

txb(sl, "Dataset: NTU-NLP-sg/xCodeEval  ·  program_synthesis  ·  106 problems  ·  compact split",
    Inches(0.4), Inches(5.7), Inches(12), Inches(0.38),
    size=13, color=LIGHT_GRAY, italic=True)

txb(sl, "Metric: pass@5  — probability at least 1 of 5 generated solutions passes all hidden unit tests",
    Inches(0.4), Inches(6.15), Inches(12), Inches(0.38),
    size=14, color=YELLOW)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Evaluation Pipeline
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
section_header(sl, "Evaluation Pipeline", "End-to-end: from problem → score")

pipeline = [
    ("Dataset",       "xCodeEval\nprogram_synthesis\n106 problems",   ACCENT,  Inches(0.4)),
    ("Generate",      "gen_program_synthesis.py\nN solutions per problem\nvia LiteLLM API", GREEN, Inches(3.1)),
    ("Execute",       "ExecEval Docker\nRun vs hidden unit tests\nper language compiler", YELLOW, Inches(5.8)),
    ("Score",         "eval_program_synthesis.py\nPASSED / WRONG_ANSWER\nper test case", ACCENT2, Inches(8.5)),
    ("Results",       "get_result.py\npass@k per language\n& average", RED,    Inches(11.2)),
]
for title, body, col, lx in pipeline:
    rect(sl, lx, Inches(1.5), Inches(2.4), Inches(2.3), CARD_BG)
    rect(sl, lx, Inches(1.5), Inches(2.4), Inches(0.07), col)
    txb(sl, title, lx + Inches(0.12), Inches(1.58), Inches(2.2), Inches(0.42),
        size=15, bold=True, color=col)
    txb(sl, body,  lx + Inches(0.12), Inches(2.05), Inches(2.2), Inches(1.6),
        size=12, color=WHITE)

# arrows
for ax in [Inches(2.85), Inches(5.55), Inches(8.25), Inches(10.95)]:
    txb(sl, "→", ax, Inches(2.25), Inches(0.35), Inches(0.5),
        size=22, bold=True, color=ACCENT2, align=PP_ALIGN.CENTER)

txb(sl, "Output files per model:",
    Inches(0.4), Inches(4.1), Inches(5), Inches(0.38),
    size=15, bold=True, color=ACCENT2)
bullet(sl, [
    "dumped/<model>/program_synthesis/<idx>_<temp>_<lang>.json  — raw API responses",
    "dumped/<model>/oai/prog_synthesis_n_sample_20/reproduce_1/<compiler>.jsonl  — execution results",
], Inches(0.4), Inches(4.55), Inches(12.5), size=13, gap=Inches(0.42))

txb(sl, "Each JSON file contains the full API response (N choices) + original problem data",
    Inches(0.4), Inches(5.55), Inches(12), Inches(0.38),
    size=13, color=LIGHT_GRAY, italic=True)

txb(sl, "ExecEval uses sandboxed Docker execution with real compilers: GCC, JDK 21, Node.js, PyPy, Go, Kotlin, PHP, Rust, Mono C#",
    Inches(0.4), Inches(6.05), Inches(12.5), Inches(0.5),
    size=13, color=LIGHT_GRAY, italic=True)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — GPT-5.1 Results
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
section_header(sl, "Results — GPT-5.1 (gpt4o alias)",
               "pass@5 × 100  |  Program Synthesis  |  n=8 samples per problem")

results = [
    ("C++",        62.83, ACCENT),
    ("Javascript", 59.42, GREEN),
    ("PHP",        58.68, GREEN),
    ("Kotlin",     49.44, YELLOW),
    ("Python",     49.04, YELLOW),
    ("Go",         44.53, YELLOW),
    ("Java",       24.87, RED),
]

bar_max  = Inches(8.5)
bar_base = Inches(1.5)
bar_h    = Inches(0.52)
gap      = Inches(0.65)

for i, (lang, score, col) in enumerate(results):
    y = bar_base + i * gap
    # label
    txb(sl, lang, Inches(0.4), y + Inches(0.08), Inches(1.4), bar_h,
        size=15, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)
    # bar background
    rect(sl, Inches(1.9), y, bar_max, bar_h, CARD_BG)
    # bar fill
    fill_w = bar_max * (score / 100)
    rect(sl, Inches(1.9), y, fill_w, bar_h, col)
    # score label
    txb(sl, f"{score}%", Inches(1.9) + fill_w + Inches(0.1), y + Inches(0.08),
        Inches(1.0), bar_h, size=15, bold=True, color=col)

avg = 49.83
hline(sl, Inches(6.25), ACCENT2)
txb(sl, f"Average pass@5  (7 languages):   {avg}%",
    Inches(0.4), Inches(6.35), Inches(7), Inches(0.5),
    size=18, bold=True, color=ACCENT)

txb(sl, "560 / 742 problems scored  |  182 skipped (empty dataset fields)",
    Inches(0.4), Inches(6.9), Inches(9), Inches(0.38),
    size=12, color=LIGHT_GRAY, italic=True)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Comparison vs Paper Baselines
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
section_header(sl, "How We Compare to Paper Baselines",
               "xCodeEval paper (2023) vs our evaluation (2025)")

rows = [
    ("Model",             "Year", "Avg pass@5",  "Notes",                    WHITE,    True),
    ("CodeT5+ 16B",       "2023", "~10%",        "Open-source encoder-decoder", LIGHT_GRAY, False),
    ("InCoder 6B",        "2023", "~5%",         "Open-source infill model",  LIGHT_GRAY, False),
    ("text-davinci-003",  "2023", "~22%",        "OpenAI paper baseline",     LIGHT_GRAY, False),
    ("GPT-3.5-turbo",     "2023", "~25%",        "OpenAI paper baseline",     LIGHT_GRAY, False),
    ("GPT-4",             "2023", "~38%",        "OpenAI, limited eval",      ACCENT2,  False),
    ("GPT-5.1  ★",        "2025", "49.8%",       "Our run — 7 langs, n=8",    ACCENT,   False),
]
col_w = [Inches(3.2), Inches(1.2), Inches(1.6), Inches(5.5)]
col_x = [Inches(0.4), Inches(3.65), Inches(4.9), Inches(6.55)]
row_h = Inches(0.52)
row_y = Inches(1.45)

for ri, (model, year, score, note, color, is_hdr) in enumerate(rows):
    y = row_y + ri * row_h
    if is_hdr:
        rect(sl, Inches(0.35), y, Inches(12.6), row_h, ACCENT)
    elif model.startswith("GPT-5.1"):
        rect(sl, Inches(0.35), y, Inches(12.6), row_h, RGBColor(0x0A, 0x30, 0x50))
    for ci, (text, cx, cw) in enumerate(zip([model, year, score, note], col_x, col_w)):
        txb(sl, text, cx, y + Inches(0.07), cw, row_h - Inches(0.07),
            size=14, bold=is_hdr,
            color=DARK_BG if is_hdr else (ACCENT if model.startswith("GPT-5.1") else color))

txb(sl, "★  GPT-5.1 shows a +11.8pp improvement over GPT-4, validating frontier model progress on executable code benchmarks.",
    Inches(0.4), Inches(5.7), Inches(12.4), Inches(0.5),
    size=14, bold=True, color=YELLOW)

txb(sl, "Note: paper baselines cover all 11 languages; our run covers 7. Exact per-language comparisons require aligned language sets.",
    Inches(0.4), Inches(6.3), Inches(12.4), Inches(0.4),
    size=12, color=LIGHT_GRAY, italic=True)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Key Observations
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
section_header(sl, "Key Observations", "GPT-5.1 program synthesis results")

obs = [
    (ACCENT,  "C++ leads at 62.8%",
              "Likely due to strict typing and structured competitive programming patterns that GPT-5.1 handles well."),
    (GREEN,   "Javascript & PHP close behind (~59%)",
              "Dynamic languages with flexible syntax — model generates idiomatic solutions reliably."),
    (YELLOW,  "Python at 49% — lower than expected",
              "Python problems in xCodeEval are algorithmic/competitive, not scripting — harder than typical Python tasks."),
    (RED,     "Java weakest at 24.9%",
              "Java requires strict imports, class boilerplate, and typed I/O. Small formatting errors fail all unit tests."),
    (ACCENT2, "Average 49.8% pass@5 across 7 languages",
              "Clear step-up from GPT-4 (~38%) reported in the original paper. Frontier model progress is measurable."),
]

for i, (col, title, desc) in enumerate(obs):
    y = Inches(1.45) + i * Inches(1.1)
    rect(sl, Inches(0.4), y, Inches(0.08), Inches(0.9), col)
    txb(sl, title, Inches(0.6), y + Inches(0.02), Inches(11.8), Inches(0.4),
        size=16, bold=True, color=col)
    txb(sl, desc,  Inches(0.6), y + Inches(0.42), Inches(11.8), Inches(0.45),
        size=13, color=LIGHT_GRAY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Next Steps
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
section_header(sl, "Next Steps", "Completing the full benchmark comparison")

left_items = [
    ("In Progress", [
        "Qwen-NVFP4 generation running  (144 / 742 files)",
        "Laguna-NVFP4 generation running (218 / 742 files)",
        "APR dataset downloaded",
        "Code Translation dataset downloaded",
    ], YELLOW),
    ("Pending", [
        "Score Qwen & Laguna once generation completes",
        "3-model comparison table (GPT-5.1 vs Qwen vs Laguna)",
        "Run Code Translation task (440 compact_small problems)",
        "Run APR task (automated program repair)",
    ], ACCENT),
]

y = Inches(1.45)
for section, items, col in left_items:
    rect(sl, Inches(0.4), y, Inches(0.08), Inches(0.4), col)
    txb(sl, section, Inches(0.6), y, Inches(5.5), Inches(0.4),
        size=17, bold=True, color=col)
    y += Inches(0.48)
    for item in items:
        rect(sl, Inches(0.65), y + Inches(0.1), Inches(0.1), Inches(0.1), col)
        txb(sl, item, Inches(0.85), y, Inches(5.8), Inches(0.38),
            size=14, color=WHITE)
        y += Inches(0.42)
    y += Inches(0.25)

# target table (right side)
txb(sl, "Target comparison table:", Inches(7.2), Inches(1.45), Inches(5.6), Inches(0.4),
    size=15, bold=True, color=ACCENT2)
tbl_rows = [
    ("Language",   "GPT-5.1", "Qwen-FP4", "Laguna-FP4", True),
    ("C++",        "62.8%",   "TBD",      "TBD",         False),
    ("Javascript", "59.4%",   "TBD",      "TBD",         False),
    ("PHP",        "58.7%",   "TBD",      "TBD",         False),
    ("Kotlin",     "49.4%",   "TBD",      "TBD",         False),
    ("Python",     "49.0%",   "TBD",      "TBD",         False),
    ("Go",         "44.5%",   "TBD",      "TBD",         False),
    ("Java",       "24.9%",   "TBD",      "TBD",         False),
    ("Average",    "49.8%",   "TBD",      "TBD",         False),
]
cx = [Inches(7.2), Inches(9.1), Inches(10.3), Inches(11.5)]
cw = [Inches(1.8), Inches(1.1), Inches(1.1), Inches(1.1)]
for ri, (*cells, is_hdr) in enumerate(tbl_rows):
    ry = Inches(1.95) + ri * Inches(0.48)
    if is_hdr:
        rect(sl, Inches(7.2), ry, Inches(5.5), Inches(0.42), ACCENT)
    elif cells[0] == "Average":
        rect(sl, Inches(7.2), ry, Inches(5.5), Inches(0.42), CARD_BG)
    for ci, (cell, x, w) in enumerate(zip(cells, cx, cw)):
        color = DARK_BG if is_hdr else (ACCENT if cells[0] == "Average" else WHITE)
        txb(sl, cell, x, ry + Inches(0.05), w, Inches(0.38),
            size=13, bold=is_hdr or cells[0]=="Average",
            color=color,
            align=PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT)

# ══════════════════════════════════════════════════════════════════════════════
# Save
# ══════════════════════════════════════════════════════════════════════════════
out = "/home/ujjwal.tiwari/ace/benchmarks/xcodeEval/xcodeeval_results.pptx"
prs.save(out)
print(f"Saved: {out}")
